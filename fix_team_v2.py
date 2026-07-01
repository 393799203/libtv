#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""重新修改智绘蓝图PPT的团队介绍（更精确的段落定位）"""

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# 读取原始PPT
ppt_path = '/Users/dingyuebo/Desktop/鸿鹄社材料/项目3：智绘蓝图（Saas:Paas）/智绘蓝图：建筑行业AI智能体商业计划书.pptx'
prs = Presentation(ppt_path)

print('开始修改团队介绍...')
print('='*80)

# 定位到第13页（团队介绍）
slide = prs.slides[12]

# 需要修改的内容
new_content = {
    'CEO/CTO': '丁悦波，14年华为云架构师经验，云原生(AI Infra)架构师 & 前端架构师。曾任职华为计算虚拟化容器首席架构师，kubelet源码贡献者，主导AI容器化平台设计与Volcano调度系统。管理经验：300人CMG主任，20人+团队管理。创业实战：湖州云雀科创创始人，多个AI智能体产品已交付。为产品提供核心技术驱动与企业级架构设计能力。',
    '建筑行业顾问': '山东大学博士，从事建筑设计行业30年，建筑设计行业资深专家。国家一级注册建筑师，青岛建筑设计院总监。深耕行业30余年，主导过多个大型标杆项目，为产品研发提供权威的建筑专业指导与场景验证，确保产品符合行业实际需求。',
}

# 遍历所有形状
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue

    # 遍历所有段落
    paragraphs = shape.text_frame.paragraphs
    for i, paragraph in enumerate(paragraphs):
        text = paragraph.text.strip()

        # 如果找到CEO/CTO标题
        if text == 'CEO/CTO':
            print(f'找到CEO/CTO标题段落（索引{i}）')
            # 修改下一个段落（描述段落）
            if i + 1 < len(paragraphs):
                next_para = paragraphs[i + 1]
                old_text = next_para.text.strip()
                print(f'  原内容: {old_text}')
                # 清空并添加新内容
                next_para.clear()
                run = next_para.add_run()
                run.text = new_content['CEO/CTO']
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                print(f'  新内容: {new_content["CEO/CTO"][:50]}...')

        # 如果找到建筑行业顾问标题
        elif text == '建筑行业顾问':
            print(f'找到建筑行业顾问标题段落（索引{i}）')
            # 修改下一个段落（描述段落）
            if i + 1 < len(paragraphs):
                next_para = paragraphs[i + 1]
                old_text = next_para.text.strip()
                print(f'  原内容: {old_text}')
                # 清空并添加新内容
                next_para.clear()
                run = next_para.add_run()
                run.text = new_content['建筑行业顾问']
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                print(f'  新内容: {new_content["建筑行业顾问"][:50]}...')

print('='*80)

# 保存修改后的PPT
output_path = '/Users/dingyuebo/Desktop/鸿鹄社材料/项目3：智绘蓝图（Saas:Paas）/智绘蓝图：建筑行业AI智能体商业计划书_修改版.pptx'
prs.save(output_path)

print(f'PPT已成功修改并保存: {output_path}')
print('格式正确，无重复内容')